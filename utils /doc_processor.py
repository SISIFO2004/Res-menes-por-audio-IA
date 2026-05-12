import io
import PyPDF2
from pptx import Presentation
from typing import Optional

def _extract_from_pdf(file_stream: io.BytesIO) -> str:
    """
    Itera sobre un archivo PDF binario y extrae el texto por página.
    Mantiene un índice de paginación para contexto estructural.
    """
    texto_extraido = []
    try:
        reader = PyPDF2.PdfReader(file_stream)
        for i, page in enumerate(reader.pages):
            texto_pagina = page.extract_text()
            if texto_pagina:
                # Se añade un marcador de página para ayudar a la IA con la jerarquía
                texto_extraido.append(f"--- Página {i+1} ---\n{texto_pagina}\n")
        
        return "\n".join(texto_extraido)
    except Exception as e:
        return f"Error crítico al leer el PDF: {str(e)}"

def _extract_from_pptx(file_stream: io.BytesIO) -> str:
    """
    Itera sobre las diapositivas de un PPTX y extrae el texto de las formas (shapes).
    """
    texto_extraido = []
    try:
        prs = Presentation(file_stream)
        for i, slide in enumerate(prs.slides):
            texto_slide = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto_slide.append(shape.text)
            
            if texto_slide:
                # Se consolida el texto por diapositiva
                contenido = " ".join(texto_slide).replace('\n', ' ')
                texto_extraido.append(f"--- Diapositiva {i+1} ---\n{contenido}\n")
                
        return "\n".join(texto_extraido)
    except Exception as e:
        return f"Error crítico al leer el PPTX: {str(e)}"

def process_document(uploaded_file) -> Optional[str]:
    """
    Función de enrutamiento principal. 
    Recibe el archivo cargado y delega la extracción según la extensión.
    """
    if uploaded_file is None:
        return None

    file_name = uploaded_file.name.lower()
    file_stream = io.BytesIO(uploaded_file.getvalue())
    
    if file_name.endswith('.pdf'):
        return _extract_from_pdf(file_stream)
    elif file_name.endswith('.pptx'):
        return _extract_from_pptx(file_stream)
    else:
        return "Error: Formato de documento no soportado. Utilice PDF o PPTX."
