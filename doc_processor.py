import PyPDF2
from pptx import Presentation
import io

def process_documents(doc_files):
    """
    Procesa una lista de archivos PDF y PPTX subidos a través de Streamlit.
    Extrae el texto plano de manera eficiente y lo estructura con etiquetas 
    de origen para facilitar la validación cruzada del LLM.
    """
    if not doc_files:
        return None
        
    # Usamos una lista para almacenar fragmentos (mucho más rápido que += en strings largos)
    text_chunks = []
    
    for doc in doc_files:
        nombre_archivo = doc.name
        text_chunks.append(f"\n\n=========== [INICIO DE DOCUMENTO: {nombre_archivo}] ===========\n")
        
        try:
            # Procesamiento de PDFs
            if nombre_archivo.lower().endswith('.pdf'):
                pdf_reader = PyPDF2.PdfReader(doc)
                for page_num, page in enumerate(pdf_reader.pages):
                    texto_pagina = page.extract_text()
                    if texto_pagina:
                        # Añadimos indicadores de página para mayor precisión del LLM
                        text_chunks.append(f"\n[Pág. {page_num + 1}] {texto_pagina.strip()}")
                        
            # Procesamiento de Diapositivas (PowerPoint)
            elif nombre_archivo.lower().endswith('.pptx'):
                # Streamlit pasa el archivo en memoria, usamos BytesIO para que python-pptx lo lea correctamente
                ppt = Presentation(io.BytesIO(doc.read()))
                for slide_num, slide in enumerate(ppt.slides):
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())
                    
                    if slide_text:
                        texto_unido = " | ".join(slide_text)
                        text_chunks.append(f"\n[Diapositiva {slide_num + 1}] {texto_unido}")
            
            else:
                text_chunks.append(f"\n[Advertencia: Formato de archivo no soportado para extracción de texto]")

        except Exception as e:
            text_chunks.append(f"\n[ERROR: Falla técnica al extraer datos de {nombre_archivo}. Detalle: {str(e)}]")
            
        text_chunks.append(f"\n=========== [FIN DE DOCUMENTO: {nombre_archivo}] ===========\n")
            
    # Unimos todos los fragmentos eficientemente y retornamos el super-string
    return "".join(text_chunks)
